"""Generate the AI Marketing Copilot presentation (PowerPoint).

Produces a 16:9, 12-slide deck mirroring video_script.md:
  1. Title  2. Problem  3. Pain Points  4. Solution Overview
  5. Demo divider  6. Setup & Seed  7. Data Sub-Agents  8. Copywriter
  9. Reviewer  10. Telegram Delivery  11. Business Impact  12. Thank You

Run:  uv run python make_presentation.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0F, 0x2A, 0x43)
NAVY_SOFT = RGBColor(0x16, 0x3A, 0x5C)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
TEAL_SOFT = RGBColor(0xE1, 0xF0, 0xF3)
AMBER = RGBColor(0xF1, 0x8F, 0x01)
AMBER_SOFT = RGBColor(0xFD, 0xF0, 0xDC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x22, 0x2B, 0x33)
MUTED = RGBColor(0x5A, 0x6B, 0x7B)
CARD_BORDER = RGBColor(0xD5, 0xDE, 0xE6)
CODE_BG = RGBColor(0x11, 0x22, 0x33)
CODE_FG = RGBColor(0xE8, 0xEE, 0xF4)
CYAN_LT = RGBColor(0x7F, 0xD4, 0xE0)
SLATE_LT = RGBColor(0xC9, 0xD8, 0xE8)
SLATE_MD = RGBColor(0x9F, 0xB4, 0xC8)

FONT = "Calibri"
FONT_MONO = "Consolas"

# Icons (kept as Unicode escapes so this file stays ASCII-only)
ICON_TIME = "\U000023F1"
ICON_BRAIN = "\U0001F9E0"
ICON_DOWN = "\U0001F4C9"
ICON_SEARCH = "\U0001F50E"
ICON_BOX = "\U0001F4E6"
ICON_STAR = "\U00002B50"
ICON_CLAP = "\U0001F3AC"
ICON_CAM = "\U0001F4F8"
ICON_CHAT = "\U0001F4AC"
ICON_IDEA = "\U0001F4A1"
ICON_SHIELD = "\U0001F6E1"
ICON_COMPASS = "\U0001F9ED"
ICON_BELL = "\U0001F514"
ICON_SUN = "\U00002600"
ICON_DUSK = "\U0001F306"
ICON_MOON = "\U0001F319"
ICON_CHECK = "\u2713"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, radius=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(st, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for t, size, color, bold in para:
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def shape_text(shp, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
               space_after=4, line_spacing=1.0):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for t, size, color, bold in para:
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold


def set_notes(s, note_text):
    s.notes_slide.notes_text_frame.text = note_text


def header(s, kicker, title, subtitle=None):
    rect(s, 0, 0, SLIDE_W, Inches(1.35), WHITE, line=CARD_BORDER, radius=False)
    rect(s, 0, 0, Inches(0.16), Inches(1.35), TEAL, radius=False)
    text(s, Inches(0.55), Inches(0.18), Inches(12.3), Inches(0.3),
         [[(kicker.upper(), 12, TEAL, True)]])
    text(s, Inches(0.55), Inches(0.44), Inches(12.3), Inches(0.55),
         [[(title, 28, NAVY, True)]])
    if subtitle:
        text(s, Inches(0.55), Inches(0.98), Inches(12.3), Inches(0.3),
             [[(subtitle, 13, MUTED, False)]])


def footer(s, idx):
    text(s, Inches(0.55), Inches(7.12), Inches(8), Inches(0.3),
         [[("AI Marketing Copilot - Presentation", 10, MUTED, False)]])
    text(s, Inches(12.4), Inches(7.12), Inches(0.6), Inches(0.3),
         [[(str(idx), 10, MUTED, True)]], align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, fill=WHITE, border=CARD_BORDER, accent=None):
    shp = rect(s, x, y, w, h, fill, line=border)
    if accent is not None:
        bar = rect(s, x, y, Inches(0.09), h, accent, radius=False)
        bar.line.fill.background()
    return shp


def full_bleed(s, color):
    rect(s, 0, 0, SLIDE_W, SLIDE_H, color, radius=False)


# ---------------------------------------------------------------------------
# Slide 1 - Title
# ---------------------------------------------------------------------------
s = slide()
full_bleed(s, NAVY)
rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), TEAL, radius=False)
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
     [[("AGENT-COPY-WRITER  /  MVP DEMO", 15, CYAN_LT, True)]])
text(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.6),
     [[("AI Marketing Copilot", 54, WHITE, True)]])
text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.6),
     [[("Turn transactional data into ready-to-publish promo content,", 20, SLATE_LT, False)]])
text(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.6),
     [[("delivered straight to the seller's Telegram.", 20, SLATE_LT, False)]])
chips = ["Python 3.14", "Google ADK (LiteLlm)", "PostgreSQL 18", "Telegram Bot"]
cx = Inches(0.9)
for c in chips:
    w = Inches(0.45 + 0.155 * len(c))
    chip = rect(s, cx, Inches(5.4), w, Inches(0.52), NAVY_SOFT, line=TEAL)
    shape_text(chip, [[(c, 13, WHITE, True)]], align=PP_ALIGN.CENTER)
    cx += w + Inches(0.25)
text(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.4),
     [[("A multi-agent solution for Shopee / TikTok Shop / Instagram Business", 14, SLATE_MD, False)]])
set_notes(s, "Good morning, and thank you for your time. Today I would like to present our AI Marketing Copilot - a multi-agent solution designed for Indonesian online sellers.")

# ---------------------------------------------------------------------------
# Slide 2 - The Problem
# ---------------------------------------------------------------------------
s = slide()
header(s, "Introduction", "The Problem: Sellers Are Overloaded",
       "A familiar situation for Indonesian online sellers across multiple platforms")
c = card(s, Inches(0.55), Inches(1.75), Inches(6.1), Inches(3.1), fill=TEAL_SOFT, border=TEAL)
shape_text(c, [
    [("2-4 hours", 54, TEAL, True)],
    [("lost every single day", 20, NAVY, True)],
    [("on repetitive operational tasks, not on growing the business.", 15, INK, False)],
], align=PP_ALIGN.CENTER, space_after=6)
text(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(4.2), [
    [("Every day, sellers manually juggle:", 16, NAVY, True)],
    [("   Managing orders and replying to customer chats", 15, INK, False)],
    [("   Monitoring inventory and finding slow-moving stock", 15, INK, False)],
    [("   Writing promotional copy for each channel from scratch", 15, INK, False)],
    [("", 8, INK, False)],
    [("The result is lost productivity, missed sales opportunities, and persistent creative block.", 15, INK, False)],
], space_after=10)
rect(s, Inches(0.55), Inches(5.2), Inches(12.25), Inches(1.0), NAVY, radius=False)
text(s, Inches(0.9), Inches(5.42), Inches(11.6), Inches(0.4),
     [[("WHO FEELS THIS PAIN?", 12, CYAN_LT, True)]])
text(s, Inches(0.9), Inches(5.74), Inches(11.6), Inches(0.4),
     [[("Shopee sellers    /    TikTok Shop sellers    /    Instagram businesses", 17, WHITE, True)]])
footer(s, 2)
set_notes(s, "Every day these sellers spend two to four hours managing orders, replying to chats, monitoring inventory, and writing promotional content across Shopee, TikTok Shop, and Instagram. This workload consumes their time and limits productivity.")

# ---------------------------------------------------------------------------
# Slide 3 - Pain Points
# ---------------------------------------------------------------------------
s = slide()
header(s, "Introduction", "Three Pain Points, One Root Cause")
cards = [
    (ICON_TIME, "Time Drain", "Hours spent daily on routine tasks - orders, chats, and inventory updates - across several disconnected platforms."),
    (ICON_BRAIN, "Creative Block", "Brainstorming fresh, channel-specific copy for every product and every platform quickly becomes exhausting."),
    (ICON_DOWN, "Missed Opportunities", "Deadstock sits forgotten and promotional moments pass because there is simply no time to act on them."),
]
cw = Inches(3.86)
gap = Inches(0.33)
x = Inches(0.55)
for icon, title, body in cards:
    c = card(s, x, Inches(1.85), cw, Inches(3.6), fill=WHITE, border=CARD_BORDER, accent=TEAL)
    shape_text(c, [
        [(icon, 30, AMBER, False)],
        [(title, 20, NAVY, True)],
        [(body, 14, INK, False)],
    ], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=10)
    x += cw + gap
bar = rect(s, Inches(0.55), Inches(5.85), Inches(12.25), Inches(0.95), TEAL_SOFT, line=TEAL)
shape_text(bar, [[("These are automation problems - and automation problems are solvable with AI.", 16, NAVY, True)]], align=PP_ALIGN.CENTER)
footer(s, 3)
set_notes(s, "Three pain points share one root cause: manual, repetitive work. Time drain, creative block, and missed opportunities all stem from sellers doing everything by hand.")

# ---------------------------------------------------------------------------
# Slide 4 - Solution Overview
# ---------------------------------------------------------------------------
s = slide()
header(s, "Solution", "A Multi-Agent Pipeline That Writes For You",
       "Data, Copywriter, Reviewer, Telegram - fully automated")
stages = [
    ("1", "Data Sub-Agents", "3 SQL queries gather accurate numbers and real reviews.", TEAL),
    ("2", "Copywriter", "One LLM call drafts 3 channel-specific formats.", TEAL),
    ("3", "Reviewer", "Scores and revises to block hallucinated claims.", TEAL),
    ("4", "Delivery", "Clean, copy-paste content sent to Telegram.", AMBER),
]
sw = Inches(2.86)
sgap = Inches(0.27)
x = Inches(0.55)
for num, title, body, accent in stages:
    c = card(s, x, Inches(1.85), sw, Inches(2.6), fill=WHITE, border=CARD_BORDER, accent=accent)
    shape_text(c, [
        [(num, 30, accent, True)],
        [(title, 17, NAVY, True)],
        [(body, 12.5, INK, False)],
    ], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=8)
    x += sw + sgap
ay = Inches(3.05)
for i in range(3):
    ax = Inches(0.55) + (sw + sgap) * (i + 1) - sgap / 2 - Inches(0.12)
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(ax), int(ay), Inches(0.24), Inches(0.3))
    ar.fill.solid()
    ar.fill.fore_color.rgb = MUTED
    ar.line.fill.background()
    ar.shadow.inherit = False
text(s, Inches(0.55), Inches(4.95), Inches(12.25), Inches(1.4), [
    [("The core idea:  ", 15, NAVY, True),
     ("SQL sub-agents provide accurate figures, the copywriter adds creative structure, "
      "the reviewer blocks false claims, and Telegram enables fast action.", 15, INK, False)],
])
footer(s, 4)
set_notes(s, "Our solution automates this entire process with a four-stage pipeline: data sub-agents, a copywriter, a reviewer, and Telegram delivery.")

# ---------------------------------------------------------------------------
# Slide 5 - Demo divider
# ---------------------------------------------------------------------------
s = slide()
full_bleed(s, NAVY)
rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), TEAL, radius=False)
text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.5),
     [[("PART 2", 16, CYAN_LT, True)]])
text(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(1.0),
     [[("Live Demonstration", 48, WHITE, True)]])
text(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.8),
     [[("From a terminal command to a ready-to-publish Telegram message - in seconds.", 18, SLATE_LT, False)]])
footer(s, 5)
set_notes(s, "Let me now demonstrate how the system works in practice, from a single terminal command to a ready-to-publish Telegram message.")

# ---------------------------------------------------------------------------
# Slide 6 - Setup & Seed
# ---------------------------------------------------------------------------
s = slide()
header(s, "Live Demonstration - 1 of 5", "Setup & Seed Data",
       "Start the database, create the schema, and load a deterministic demo dataset")
term = rect(s, Inches(0.55), Inches(1.85), Inches(6.3), Inches(3.4), CODE_BG, radius=False)
tf = term.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.margin_right = Inches(0.2)
lines = [
    ("$ ", "docker compose up -d"),
    ("$ ", "uv run python -m app.cli init-db"),
    ("$ ", "uv run python -m app.cli seed"),
]
for i, (prefix, cmd) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    r0 = p.add_run(); r0.text = prefix
    r0.font.name = FONT_MONO; r0.font.size = Pt(16); r0.font.color.rgb = AMBER; r0.font.bold = True
    r1 = p.add_run(); r1.text = cmd
    r1.font.name = FONT_MONO; r1.font.size = Pt(16); r1.font.color.rgb = CODE_FG
text(s, Inches(7.3), Inches(1.9), Inches(5.5), Inches(0.4),
     [[("Demo dataset (fully deterministic):", 15, NAVY, True)]])
facts = [("7", "fashion SKUs"), ("94", "orders"), ("23", "customer reviews")]
fy = Inches(2.45)
for num, label in facts:
    c = card(s, Inches(7.3), fy, Inches(5.5), Inches(0.85), fill=TEAL_SOFT, border=TEAL)
    shape_text(c, [[(num + "  ", 20, TEAL, True), (label, 15, INK, False)]], align=PP_ALIGN.LEFT)
    fy += Inches(1.05)
text(s, Inches(0.55), Inches(5.6), Inches(12.25), Inches(0.9), [
    [("Every run is reproducible on any machine, so the demo behaves identically for every audience.", 14, MUTED, False)],
])
footer(s, 6)
set_notes(s, "We start the PostgreSQL database, initialise the schema, and load a deterministic demo dataset: seven products, ninety-four orders, and twenty-three reviews.")

# ---------------------------------------------------------------------------
# Slide 7 - Data Sub-Agents
# ---------------------------------------------------------------------------
s = slide()
header(s, "Live Demonstration - 2 of 5", "Three Data Sub-Agents",
       "Accurate SQL queries - not LLM guesses - assemble the driving data")
cards = [
    (ICON_SEARCH, "Sales Velocity Monitor", "Best-selling SKUs in the last 24 hours and 7 days, by orders and revenue - for FOMO-driven promos."),
    (ICON_BOX, "Deadstock Monitor", "SKUs piling up (40+ pcs, listed over 30 days, low 7-day orders) - clearance-sale candidates."),
    (ICON_STAR, "Social Proof Miner", "The 3 most recent 5-star reviews plus the average rating - authentic testimonials for the copy."),
]
cw = Inches(3.86)
gap = Inches(0.33)
x = Inches(0.55)
for icon, title, body in cards:
    c = card(s, x, Inches(1.85), cw, Inches(3.3), fill=WHITE, border=CARD_BORDER, accent=TEAL)
    shape_text(c, [
        [(icon, 28, AMBER, False)],
        [(title, 17, NAVY, True)],
        [(body, 13, INK, False)],
    ], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=10)
    x += cw + gap
text(s, Inches(0.55), Inches(5.5), Inches(12.25), Inches(0.9), [
    [("They run in parallel and return a single context object - the only source of figures the copywriter may use.", 14, MUTED, False)],
])
footer(s, 7)
set_notes(s, "Three data sub-agents query the database in parallel. The Sales Velocity Monitor finds best-sellers, the Deadstock Monitor detects slow-moving stock, and the Social Proof Miner collects recent five-star reviews.")

# ---------------------------------------------------------------------------
# Slide 8 - Copywriter
# ---------------------------------------------------------------------------
s = slide()
header(s, "Live Demonstration - 3 of 5", "One Prompt, Three Ready-to-Use Formats",
       "The ADK Copywriter Agent drafts every channel's content at once")
formats = [
    (ICON_CLAP, "TikTok & Shopee Video", "15-second script with [Visual], [Hook], [Body], and [CTA] markers."),
    (ICON_CAM, "Instagram & Shopee Feed", "Testimonial-led storytelling, emoji bullets, CTA, and hashtags."),
    (ICON_CHAT, "Shopee Broadcast / WhatsApp", "Under 50 words, urgent, direct-to-link, no hashtags."),
]
cw = Inches(3.86)
gap = Inches(0.33)
x = Inches(0.55)
for icon, title, body in formats:
    c = card(s, x, Inches(1.85), cw, Inches(2.6), fill=WHITE, border=CARD_BORDER, accent=TEAL)
    shape_text(c, [
        [(icon, 26, AMBER, False)],
        [(title, 16, NAVY, True)],
        [(body, 13, INK, False)],
    ], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=8)
    x += cw + gap
bar = rect(s, Inches(0.55), Inches(4.8), Inches(12.25), Inches(1.5), TEAL_SOFT, line=TEAL)
shape_text(bar, [
    [("Brand voice as a config file", 15, NAVY, True)],
    [("Tone, audience, CTA rules, hashtags, and forbidden claims live in one editable profile - change the file once and the voice updates everywhere.", 13.5, INK, False)],
], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, space_after=6)
footer(s, 8)
set_notes(s, "The results are assembled into one context and passed to the Copywriter Agent, which produces three channel-specific formats in a single step, guided by the brand voice profile.")

# ---------------------------------------------------------------------------
# Slide 9 - Reviewer
# ---------------------------------------------------------------------------
s = slide()
header(s, "Live Demonstration - 4 of 5", "The Reviewer Blocks Hallucinations",
       "A second agent audits every claim against the factual data before anything is sent")
text(s, Inches(0.55), Inches(1.85), Inches(5.9), Inches(0.4),
     [[("What it verifies:", 15, NAVY, True)]])
checks = [
    "Hook strength (is the 3-second hook compelling?)",
    "Figure accuracy (stock, discount, rating, orders)",
    "Brand tone and channel-appropriate CTA",
    "Format compliance for all three content types",
]
cy = Inches(2.3)
for item in checks:
    text(s, Inches(0.75), cy, Inches(5.7), Inches(0.6),
         [[(ICON_CHECK + "  ", 14, TEAL, True), (item, 14, INK, False)]], space_after=0)
    cy += Inches(0.62)
term = rect(s, Inches(6.75), Inches(1.85), Inches(6.0), Inches(3.3), CODE_BG, radius=False)
tf = term.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.margin_right = Inches(0.2)
out = [
    ICON_COMPASS + " Target: Premium Linen Shirt (45 pcs, 30% off)",
    "",
    "Review 1/2: score=78  approved=False",
    "  -> feed needs emoji bullets, stronger CTA",
    "",
    "Review 2/2: score=96  approved=True",
    "",
    "[OK] Approved=True | rounds=2",
]
for i, line in enumerate(out):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(3)
    r = p.add_run(); r.text = line
    r.font.name = FONT_MONO; r.font.size = Pt(13); r.font.color.rgb = CODE_FG
text(s, Inches(0.55), Inches(5.6), Inches(12.25), Inches(0.8), [
    [("Reject, revise, re-review - up to 2 rounds. Only approved content proceeds to delivery.", 14, MUTED, False)],
])
footer(s, 9)
set_notes(s, "Before delivery, the Reviewer Agent audits every figure against the source data to prevent hallucination. It scores the draft from one to one hundred and requests a revision whenever quality falls short.")

# ---------------------------------------------------------------------------
# Slide 10 - Telegram Delivery
# ---------------------------------------------------------------------------
s = slide()
header(s, "Live Demonstration - 5 of 5", "Delivered to Telegram, Ready to Publish",
       "The seller receives a clean, formatted message they can copy immediately")
phone = rect(s, Inches(0.7), Inches(1.85), Inches(4.4), Inches(4.9), NAVY, radius=True)
rect(s, Inches(1.0), Inches(2.0), Inches(3.8), Inches(0.4), TEAL, radius=False)
bubble = rect(s, Inches(1.05), Inches(2.6), Inches(3.7), Inches(3.7), WHITE, radius=True)
shape_text(bubble, [
    [(ICON_BELL + " READY-TO-USE CONTENT", 11, NAVY, True)],
    [("PRIME TIME MORNING", 11, TEAL, True)],
    [("", 6, INK, False)],
    [(ICON_CLAP + " Video script (~15s)", 10.5, NAVY, True)],
    [(ICON_CAM + " Feed caption", 10.5, NAVY, True)],
    [(ICON_CHAT + " Broadcast chat", 10.5, NAVY, True)],
    [("", 6, INK, False)],
    [("Formatted and copy-paste ready.", 10.5, MUTED, False)],
], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=5)
text(s, Inches(5.6), Inches(2.0), Inches(7.1), Inches(4.6), [
    [("One command runs the whole pipeline:", 15, NAVY, True)],
    [("", 4, INK, False)],
    [("uv run python -m app.cli trigger --slot morning", 14, AMBER, True)],
    [("", 4, INK, False)],
    [("   Real delivery via Telegram bot token + chat ID", 14, INK, False)],
    [("   --dry-run prints to terminal and archives to outputs/", 14, INK, False)],
    [("   Every run is archived with a timestamp for audit", 14, INK, False)],
    [("", 6, INK, False)],
    [("Prime-time slots: morning " + ICON_SUN + " / evening " + ICON_DUSK + " / night " + ICON_MOON, 14, NAVY, True)],
], space_after=10)
footer(s, 10)
set_notes(s, "Finally, the approved content is delivered directly to the seller's Telegram, formatted and ready to copy and publish. What previously took hours is now completed automatically in seconds.")

# ---------------------------------------------------------------------------
# Slide 11 - Business Impact
# ---------------------------------------------------------------------------
s = slide()
header(s, "Business Impact", "What This Means for Online Sellers",
       "Measurable outcomes, not just a technical novelty")
cards = [
    (ICON_TIME, "Time Saved", "2-4 hours recovered every day - sellers reclaim their working day for growth, not routine."),
    (ICON_IDEA, "Creative Block Solved", "Channel-specific content generated on demand, so sellers never stare at a blank page again."),
    (ICON_SHIELD, "Brand Trust Protected", "An automated quality loop prevents inaccurate claims and keeps the brand's voice consistent."),
]
cw = Inches(3.86)
gap = Inches(0.33)
x = Inches(0.55)
for icon, title, body in cards:
    c = card(s, x, Inches(1.85), cw, Inches(3.4), fill=WHITE, border=CARD_BORDER, accent=AMBER)
    shape_text(c, [
        [(icon, 30, AMBER, False)],
        [(title, 19, NAVY, True)],
        [(body, 13.5, INK, False)],
    ], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=10)
    x += cw + gap
text(s, Inches(0.55), Inches(5.6), Inches(12.25), Inches(0.8), [
    [("Operate more efficiently, respond faster to market demand, and grow the business sustainably.", 15, NAVY, True)],
])
footer(s, 11)
set_notes(s, "In conclusion, this solution saves sellers two to four hours daily, removes creative block, and protects brand trust through an automated quality-review loop.")

# ---------------------------------------------------------------------------
# Slide 12 - Thank You
# ---------------------------------------------------------------------------
s = slide()
full_bleed(s, NAVY)
rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), TEAL, radius=False)
text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.0),
     [[("Thank You", 54, WHITE, True)]])
text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.8),
     [[("AI Marketing Copilot - empowering Indonesian online sellers to work smarter.", 18, SLATE_LT, False)]])
text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.5),
     [[("Questions welcome.", 14, CYAN_LT, True)]])
set_notes(s, "Thank you very much for your attention. Questions are welcome.")

prs.save("ai-marketing-copilot-presentation.pptx")
print("Saved ai-marketing-copilot-presentation.pptx with", len(prs.slides._sldIdLst), "slides.")